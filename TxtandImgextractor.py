import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# =========================
# CONFIG
# =========================

folder_path = r"C:\Users\ASUS\Desktop\Exam"

output_file = os.path.join(folder_path, "combined_output.txt")
image_folder = os.path.join(folder_path, "images")

os.makedirs(image_folder, exist_ok=True)

image_counter = 1


# =========================
# PDF PROCESSING
# =========================

def process_pdf(file_path):
    global image_counter
    text = ""

    try:
        doc = fitz.open(file_path)

        for page_num in range(len(doc)):
            page = doc[page_num]



            # TEXT
            try:
                text += page.get_text()
            except Exception:
                pass

            # IMAGES
            try:
                for img in page.get_images(full=True):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)

                        image_bytes = base_image.get("image")
                        image_ext = base_image.get("ext", "png")

                        if not image_bytes:
                            continue

                        img_name = f"image_{image_counter}.{image_ext}"
                        img_path = os.path.join(image_folder, img_name)

                        with open(img_path, "wb") as f:
                            f.write(image_bytes)


                        image_counter += 1

                    except:
                        pass
            except:
                pass

    except Exception as e:
        print(f"❌ PDF Error: {file_path} -> {e}")

    return text


# =========================
# PPTX PROCESSING
# =========================

def process_pptx(file_path):
    global image_counter
    text = ""

    try:
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):

            for shape in slide.shapes:

                # TEXT
                try:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                except:
                    pass

                # IMAGE
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image

                        img_name = f"image_{image_counter}.{image.ext}"
                        img_path = os.path.join(image_folder, img_name)

                        with open(img_path, "wb") as f:
                            f.write(image.blob)


                        image_counter += 1
                except:
                    pass

    except Exception as e:
        print(f"❌ PPTX Error: {file_path} -> {e}")

    return text


# =========================
# MAIN
# =========================

def main():
    files = sorted(os.listdir(folder_path), key=str.lower)

    print("\n🔍 ALL FILES DETECTED:")
    for f in files:
        print(" -", f)

    with open(output_file, "w", encoding="utf-8") as outfile:

        for file in files:
            file_path = os.path.join(folder_path, file)

            if not os.path.isfile(file_path):
                continue

            file_lower = file.lower()

            if file_lower.endswith((".pdf", ".pptx", ".ppt")):

                print(f"\n📄 Processing: {file}")

                outfile.write("\n" + "=" * 80 + "\n")
                outfile.write(f"FILE: {file}\n")
                outfile.write("=" * 80 + "\n")

                # PDF
                if file_lower.endswith(".pdf"):
                    outfile.write(process_pdf(file_path))

                # PPTX
                elif file_lower.endswith(".pptx"):
                    outfile.write(process_pptx(file_path))

                # PPT → CONVERT + PROCESS
                elif file_lower.endswith(".ppt"):
                    print(f"⚠️ Skipping .ppt (convert to .pptx manually): {file}")

            else:
                print(f"⏭️ Skipped: {file}")

    print("\n✅ DONE!")
    print("📄 Output:", output_file)
    print("🖼️ Images:", image_folder)


# =========================
# RUN
# =========================

if __name__ == "__main__":
    main()
